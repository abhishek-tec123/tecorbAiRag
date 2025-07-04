import os
import logging
from concurrent.futures import ThreadPoolExecutor, as_completed
from typing import List, Dict, Optional
from docx import Document
import pandas as pd
from pptx import Presentation

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(message)s'
)
logger = logging.getLogger(__name__)

class TextExtractor:
    """
    Extract text from various file types: txt, pdf, docx, xlsx, xls, csv, pptx.
    """
    def __init__(self, file_path: str):
        if not os.path.isfile(file_path):
            logger.error(f"File not found: {file_path}")
            raise FileNotFoundError(f"File not found: {file_path}")
        self.file_path = file_path
        self.extension = os.path.splitext(file_path)[1].lower()

    def extract_text(self, preview: Optional[int] = None) -> str:
        """
        Extract text from the file, with optional preview length.
        """
        logger.info(f"Starting extraction for {self.file_path} (type: {self.extension})")
        try:
            if self.extension == '.txt':
                result = self._extract_txt(preview)
            elif self.extension == '.pdf':
                result = self._extract_pdf(preview)
            elif self.extension == '.docx':
                result = self._extract_docx(preview)
            elif self.extension in ['.xlsx', '.xls']:
                result = self._extract_excel(preview)
            elif self.extension == '.csv':
                result = self._extract_csv(preview)
            elif self.extension == '.pptx':
                result = self._extract_pptx(preview)
            else:
                logger.error(f"Unsupported file type: {self.extension}")
                raise ValueError(f"Unsupported file type: {self.extension}")
            logger.info(f"Extraction successful for {self.file_path}")
            return result
        except Exception as e:
            logger.error(f"Failed to extract text from {self.file_path}: {str(e)}")
            return f"[Error] Failed to extract text from {self.file_path}: {str(e)}"

    def _extract_txt(self, preview: Optional[int] = None) -> str:
        logger.debug(f"Extracting TXT: {self.file_path}")
        with open(self.file_path, 'r', encoding='utf-8') as f:
            return f.read(preview) if preview is not None else f.read()

    def _extract_pdf(self, preview: Optional[int] = None) -> str:
        logger.debug(f"Extracting PDF: {self.file_path}")
        try:
            import fitz  # PyMuPDF
            text = ""
            with fitz.open(self.file_path) as doc:
                for page in doc:
                    text += page.get_text()
                    if preview is not None and len(text) >= preview:
                        logger.debug(f"PDF preview limit reached for {self.file_path}")
                        return text[:preview]
            return text if preview is None else text[:preview]
        except ImportError:
            import PyPDF2
            text = ""
            with open(self.file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text
                    if preview is not None and len(text) >= preview:
                        logger.debug(f"PDF preview limit reached for {self.file_path}")
                        return text[:preview]
            return text if preview is None else text[:preview]

    def _extract_docx(self, preview: Optional[int] = None) -> str:
        logger.debug(f"Extracting DOCX: {self.file_path}")
        doc = Document(self.file_path)
        texts = []
        for para in doc.paragraphs:
            texts.append(para.text)
            if preview is not None and sum(len(t) for t in texts) >= preview:
                logger.debug(f"DOCX preview limit reached for {self.file_path}")
                break
        text = '\n'.join(texts)
        return text if preview is None else text[:preview]

    def _extract_excel(self, preview: Optional[int] = None) -> str:
        logger.debug(f"Extracting Excel: {self.file_path}")
        df_dict = pd.read_excel(self.file_path, sheet_name=None, nrows=preview if preview is not None else None)
        text = ""
        for sheet, df in df_dict.items():
            text += f"\n--- Sheet: {sheet} ---\n"
            text += df.astype(str).to_string(index=False)
            if preview is not None and len(text) >= preview:
                logger.debug(f"Excel preview limit reached for {self.file_path}")
                return text[:preview]
        return text if preview is None else text[:preview]

    def _extract_csv(self, preview: Optional[int] = None) -> str:
        logger.debug(f"Extracting CSV: {self.file_path}")
        df = pd.read_csv(self.file_path, nrows=preview if preview is not None else None)
        text = df.astype(str).to_string(index=False)
        return text if preview is None else text[:preview]

    def _extract_pptx(self, preview: Optional[int] = None) -> str:
        logger.debug(f"Extracting PPTX: {self.file_path}")
        prs = Presentation(self.file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
                    if preview is not None and len(text) >= preview:
                        logger.debug(f"PPTX preview limit reached for {self.file_path}")
                        return text[:preview]
        return text if preview is None else text[:preview]

def batch_extract_text(file_paths: List[str], max_workers: int = 4, preview: Optional[int] = None) -> Dict[str, str]:
    """
    Extract text from a batch of files in parallel.
    """
    logger.info(f"Starting batch extraction for {len(file_paths)} files.")
    results = {}
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        future_to_path = {
            executor.submit(TextExtractor(path).extract_text, preview): path for path in file_paths
        }
        for future in as_completed(future_to_path):
            path = future_to_path[future]
            try:
                results[path] = future.result()
                logger.info(f"Batch extraction complete for: {path}")
            except Exception as exc:
                results[path] = f"[Error] {exc}"
                logger.error(f"Batch extraction failed for {path}: {exc}")
    logger.info("Batch extraction finished.")
    return results

# # Example usage:
# if __name__ == "__main__":
#     files = ["/Users/abhishek/Desktop/ragTecorbAI/docs-pdf/tutorial.pdf"]
#     extracted = batch_extract_text(files, max_workers=8,preview=50)
#     for path, text in extracted.items():
#         print(f"--- {path} ---\n{text}\n")